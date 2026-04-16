"""Report Generator tool.

Generates preliminary PowerPoint and final PDF reports from analysis results.
"""

from pathlib import Path
from typing import Any

from loguru import logger
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)


class ReportGenerator:
    """Generates PPTX (preliminary) and PDF (final) analysis reports."""

    def __init__(self, output_dir: str = "src/data/reports") -> None:
        self._output_dir = Path(output_dir)
        self._output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"ReportGenerator initialized (output_dir={self._output_dir}).")

    def generate_pptx(
        self,
        title: str,
        sections: list[dict[str, Any]],
        output_name: str = "report.pptx",
    ) -> str:
        """Generate a preliminary PowerPoint report.

        Args:
            title: Report title.
            sections: List of dicts with keys 'heading', 'text', and optional 'image_path'.
            output_name: Output file name.

        Returns:
            Full path to the generated .pptx file.
        """
        prs = Presentation()

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Automated SIT Signal Analysis Report"

        # Content slides
        for section in sections:
            slide_layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section.get("heading", "")

            body = slide.placeholders[1]
            body.text = section.get("text", "")

            if "image_path" in section and Path(section["image_path"]).is_file():
                slide.shapes.add_picture(
                    section["image_path"],
                    Inches(1),
                    Inches(3),
                    width=Inches(8),
                )

        output_path = self._output_dir / output_name
        prs.save(str(output_path))
        logger.info(f"PPTX report generated: {output_path}")
        return str(output_path)

    def generate_pdf(
        self,
        title: str,
        sections: list[dict[str, Any]],
        output_name: str = "report.pdf",
    ) -> str:
        """Generate a final PDF report.

        Args:
            title: Report title.
            sections: List of dicts with keys 'heading', 'text', and optional 'image_path'.
            output_name: Output file name.

        Returns:
            Full path to the generated .pdf file.
        """
        output_path = self._output_dir / output_name
        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story: list[Any] = []

        # Title
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 24))

        for section in sections:
            heading = section.get("heading", "")
            text = section.get("text", "")

            if heading:
                story.append(Paragraph(heading, styles["Heading2"]))
                story.append(Spacer(1, 12))
            if text:
                story.append(Paragraph(text, styles["BodyText"]))
                story.append(Spacer(1, 8))
            if "image_path" in section and Path(section["image_path"]).is_file():
                img = Image(section["image_path"], width=400, height=250)
                story.append(img)
                story.append(Spacer(1, 12))

            story.append(PageBreak())

        doc.build(story)
        logger.info(f"PDF report generated: {output_path}")
        return str(output_path)
